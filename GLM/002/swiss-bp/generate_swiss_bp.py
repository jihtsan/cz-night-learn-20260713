# -*- coding: utf-8 -*-
"""瑞士国际平面风(Swiss / International Typographic Style) 商业计划书生成器。
内容来源:002/docs 下的 BP 逐页路演稿、内容框架、资金规划。
输出:002/PPT/商业计划书-瑞士风-天微AI实验室.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.oxml.ns import qn
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION

# ---------- 设计系统 / Design System ----------
SW, SH = 13.333, 7.5
M = 0.67                      # 页边距 margin
CW = SW - 2 * M               # 内容宽 = 11.99
TOTAL = 22

INK    = RGBColor(0x1A, 0x1A, 0x1A)
RED    = RGBColor(0xD9, 0x25, 0x1D)
BLUE   = RGBColor(0x1A, 0x4F, 0xA0)
GRAY   = RGBColor(0x66, 0x66, 0x66)
LGRAY  = RGBColor(0x99, 0x99, 0x99)
BORDER = RGBColor(0xE8, 0xE8, 0xE8)
BG2    = RGBColor(0xF4, 0xF4, 0xF4)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
SOFTRED = RGBColor(0xFB, 0xEC, 0xEC)

F_BLACK = "Arial Black"
F_SANS  = "Arial"
F_EA    = "微软雅黑"

prs = Presentation()
prs.slide_width  = Inches(SW)
prs.slide_height = Inches(SH)
BLANK = prs.slide_layouts[6]


def _set_font(run, latin=F_SANS, ea=F_EA, size=14, bold=False, color=INK, spc=None, italic=False):
    f = run.font
    f.size = Pt(size); f.bold = bold; f.italic = italic; f.name = latin
    f.color.rgb = color
    rPr = run._r.get_or_add_rPr()
    ea_el = rPr.find(qn('a:ea'))
    if ea_el is None:
        ea_el = rPr.makeelement(qn('a:ea'), {}); rPr.append(ea_el)
    ea_el.set('typeface', ea)
    if spc is not None:
        rPr.set('spc', str(spc))


def r(t, **kw):
    kw['t'] = t
    return kw


def add_para(tf, runs, first=False, align=PP_ALIGN.LEFT, sb=0, sa=0, ls=None, level=0):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align; p.level = level
    if sb: p.space_before = Pt(sb)
    if sa: p.space_after = Pt(sa)
    if ls: p.line_spacing = ls
    for rr in runs:
        run = p.add_run(); run.text = rr['t']
        _set_font(run, latin=rr.get('latin', F_SANS), ea=rr.get('ea', F_EA),
                  size=rr.get('size', 14), bold=rr.get('bold', False),
                  color=rr.get('color', INK), spc=rr.get('spc'), italic=rr.get('italic', False))
    return p


def tbox(slide, l, t, w, h, anchor=MSO_ANCHOR.TOP, wrap=True):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    return tf


def text(slide, l, t, w, h, runs, anchor=MSO_ANCHOR.TOP, align=PP_ALIGN.LEFT, ls=None, sa=0, sb=0, wrap=True):
    tf = tbox(slide, l, t, w, h, anchor=anchor, wrap=wrap)
    add_para(tf, runs, first=True, align=align, sb=sb, sa=sa, ls=ls)
    return tf


def rect(slide, l, t, w, h, fill=WHITE, line=None, weight=0.75):
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(weight)
    sp.shadow.inherit = False
    tf = sp.text_frame
    tf.margin_left = Inches(0.12); tf.margin_right = Inches(0.12)
    tf.margin_top = Inches(0.06); tf.margin_bottom = Inches(0.06)
    tf.word_wrap = True
    return sp


def hline(slide, l, t, w, color=BORDER, weight=1.0):
    cn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(l), Inches(t), Inches(l + w), Inches(t))
    cn.line.color.rgb = color; cn.line.width = Pt(weight)
    return cn


def vline(slide, l, t, h, color=BORDER, weight=1.0):
    cn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(l), Inches(t), Inches(l), Inches(t + h))
    cn.line.color.rgb = color; cn.line.width = Pt(weight)
    return cn


def page_num(slide, idx):
    text(slide, M, 0.42, 1.5, 0.3, [r("P · %02d" % idx, size=10, color=LGRAY, spc=200, latin=F_SANS)])


def page_foot(slide, idx):
    text(slide, SW - M - 2.2, 6.95, 2.2, 0.3,
         [r("%02d / %d" % (idx, TOTAL), size=10, color=LGRAY, spc=120, latin=F_SANS)],
         align=PP_ALIGN.RIGHT)


def page_head(slide, idx, title, en):
    page_num(slide, idx)
    text(slide, M, 0.92, CW, 0.8, [r(title, size=30, bold=True, color=INK, latin=F_BLACK)])
    text(slide, M, 1.62, CW, 0.3, [r(en.upper(), size=11, color=GRAY, spc=220, latin=F_SANS)])
    hline(slide, M, 1.98, 0.72, color=RED, weight=2.5)
    page_foot(slide, idx)


def bg_white(slide):
    rect(slide, 0, 0, SW, SH, fill=WHITE, line=None)


def add_slide():
    s = prs.slides.add_slide(BLANK)
    bg_white(s)
    return s


def notes(slide, body):
    slide.notes_slide.notes_text_frame.text = body


def cols(n, gap=0.5):
    w = (CW - gap * (n - 1)) / n
    return [(M + i * (w + gap), w) for i in range(n)]


def slide_01():
    s = add_slide()
    text(s, M, 0.5, 6, 0.3, [r("BUSINESS PLAN  ·  种子轮  ·  2026", size=11, color=GRAY, spc=240, latin=F_SANS)])
    hline(s, M, 0.84, 2.4, color=INK, weight=1.0)
    text(s, SW - M - 4.5, 0.5, 4.5, 0.3,
         [r("天微AI实验室  ·  TIANWEI AI LAB", size=11, color=GRAY, spc=160, latin=F_SANS)],
         align=PP_ALIGN.RIGHT)
    text(s, M, 3.15, 12, 1.3, [r("金融舆情", size=66, bold=True, color=INK, latin=F_BLACK)], ls=1.0)
    text(s, M, 4.32, 12, 1.2, [r("智能体平台", size=66, bold=True, color=INK, latin=F_BLACK)], ls=1.0)
    vline(s, M, 5.75, 0.95, color=RED, weight=3.0)
    text(s, M + 0.22, 5.78, 11, 0.4,
         [r("让机器像资深分析师一样 · 理解 / 研判 / 应对金融舆情", size=15, color=INK, latin=F_SANS, ea=F_EA)])
    text(s, M + 0.22, 6.12, 11, 0.4,
         [r("FINANCIAL SENTIMENT AI AGENT PLATFORM", size=11, color=GRAY, spc=220, latin=F_SANS)])
    text(s, M, 6.72, 8, 0.3,
         [r("种子轮  ·  拟融资人民币 1,000 万  ·  跑道 18–24 个月", size=11, color=GRAY, spc=140, latin=F_SANS, ea=F_EA)])
    text(s, SW - M - 2.2, 6.95, 2.2, 0.3,
         [r("01 / %d" % TOTAL, size=10, color=LGRAY, spc=120, latin=F_SANS)], align=PP_ALIGN.RIGHT)
    notes(s, "各位投资人好,我是天微AI实验室。今天汇报我们的金融舆情智能体平台:用 AI Agent 重建金融机构的舆情分析与声誉风险管理工作流。本轮种子轮计划融资 1000 万,用于验证产品市场匹配、打磨金融行业模型、拿下首批标杆客户。")


def slide_02():
    s = add_slide()
    page_head(s, 2, "目录", "Contents")
    items = [
        ("01", "项目情况介绍", "PROJECT OVERVIEW", "痛点与机会"),
        ("02", "产品方案分析", "PRODUCT & SOLUTION", "架构与壁垒"),
        ("03", "市场与竞争",   "MARKET & COMPETITION", "规模与定位"),
        ("04", "商业化与未来规划", "BUSINESS & FUTURE", "融资与里程碑"),
    ]
    c4 = cols(2, gap=0.6); ys = [2.55, 4.55]
    for i, (no, cn, en, sub) in enumerate(items):
        x, w = c4[i % 2]; y = ys[i // 2]
        text(s, x, y, w, 0.9, [r(no, size=54, bold=True, color=BORDER, latin=F_BLACK)])
        hline(s, x, y + 1.0, w, color=INK, weight=1.0)
        text(s, x, y + 1.15, w, 0.5, [r(cn, size=22, bold=True, color=INK, latin=F_BLACK)])
        text(s, x, y + 1.62, w, 0.3, [r(en, size=10, color=GRAY, spc=200, latin=F_SANS)])
        text(s, x, y + 1.95, w, 0.3, [r(sub, size=12, color=RED, latin=F_SANS, ea=F_EA)])
    notes(s, "汇报分四部分:先看行业痛点与机会,再讲产品方案与壁垒,然后是市场与竞争,最后落到商业化与融资规划。")


def section(idx, part_no, title, en, statement):
    s = add_slide()
    page_num(s, idx)
    text(s, M, 1.0, 6, 0.5, [r("PART", size=14, color=GRAY, spc=260, latin=F_SANS)])
    text(s, M, 1.25, 6, 2.0, [r(part_no, size=180, bold=True, color=BORDER, latin=F_BLACK)])
    text(s, M, 4.15, 11, 0.8, [r(title, size=40, bold=True, color=INK, latin=F_BLACK)])
    text(s, M, 4.95, 11, 0.35, [r(en.upper(), size=12, color=GRAY, spc=240, latin=F_SANS)])
    hline(s, M, 5.45, 0.72, color=RED, weight=3.0)
    text(s, M, 5.62, 11, 0.5, [r(statement, size=16, color=INK, latin=F_SANS, ea=F_EA)])
    page_foot(s, idx)
    return s


def slide_03():
    notes(section(3, "01", "项目情况介绍", "Project Overview", "金融舆情正从「关键词时代」迈入「智能体时代」。"),
          "进入第一部分:项目情况介绍。核心判断是——金融舆情正从关键词时代迈入智能体时代。")


def slide_08():
    notes(section(8, "02", "产品方案分析", "Product & Solution", "金融行业大模型 + 多 Agent 编排 + 金融知识图谱。"),
          "第二部分:产品方案分析。我们的技术底座是金融行业大模型、多 Agent 编排与金融知识图谱。")


def slide_13():
    notes(section(13, "03", "市场与竞争", "Market & Competition", "金融舆情与声誉风险管理市场,规模可观、生态位空白。"),
          "第三部分:市场与竞争。我们看金融舆情与声誉风险管理市场,天花板足够高,且关键生态位目前空白。")


def slide_18():
    notes(section(18, "04", "商业化与未来规划", "Business & Future", "种子轮:验证 PMF,18–24 个月跑道。"),
          "第四部分:商业化与未来规划。种子轮的核心任务是验证 PMF,18 到 24 个月跑道。")


def slide_04():
    s = add_slide()
    page_head(s, 4, "行业痛点", "Structural Pain Points")
    pts = [
        ("01", "关键词舆情不懂金融语义", "KEYWORD-BLIND", "无法区分「公告利好」与「被立案调查」,误报漏报严重。"),
        ("02", "人工研判慢、稀缺、贵", "SLOW & COSTLY", "潜在负面响应常需数十分钟到数小时,声誉窗口常只有分钟级。"),
        ("03", "多源多语言多模态碎片化", "FRAGMENTED", "社媒、新闻、股吧、研报、公告、短视频、境外媒体难以整合关联。"),
        ("04", "强监管下幻觉与留痕是红线", "COMPLIANCE-CRITICAL", "银行 / 券商 / 上市公司要求可解释、可追溯,容错率低。"),
    ]
    c = cols(2, gap=0.7); ys = [2.5, 4.62]
    for i, (no, cn, en, desc) in enumerate(pts):
        x, w = c[i % 2]; y = ys[i // 2]
        text(s, x, y, w, 0.85, [r(no, size=46, bold=True, color=BORDER, latin=F_BLACK)])
        hline(s, x, y + 0.95, w, color=INK, weight=1.0)
        text(s, x, y + 1.08, w, 0.5, [r(cn, size=18, bold=True, color=INK, latin=F_BLACK)])
        text(s, x, y + 1.52, w, 0.3, [r(en, size=9.5, color=RED, spc=200, latin=F_SANS)])
        text(s, x, y + 1.85, w, 0.6, [r(desc, size=12.5, color=GRAY, latin=F_SANS, ea=F_EA)], ls=1.2)
    notes(s, "金融舆情的核心矛盾:信息海量且专业,但现有工具要么是关键词+规则不懂金融,要么是人工慢且贵,同时还要满足严苛的合规留痕。这就是我们的切入口。")


def slide_05():
    s = add_slide()
    page_head(s, 5, "解决方案", "From Seeing to Acting")
    blocks = [
        ("01", "智能监测", "MONITORING", "7×24 全网 + 授权源,金融实体级(公司 / 人物 / 产品 / 事件)精准命中。"),
        ("02", "语义研判", "ANALYSIS", "事件类型、风险等级、传播态势、利益相关方自动识别与关联。"),
        ("03", "自动应对", "ACTION", "研判报告 + 应对口径草稿(人在回路确认),分级预警触达。"),
    ]
    c = cols(3, gap=0.6)
    for i, (no, cn, en, desc) in enumerate(blocks):
        x, w = c[i]
        text(s, x, 2.55, w, 0.9, [r(no, size=46, bold=True, color=BORDER, latin=F_BLACK)])
        hline(s, x, 3.5, w, color=INK, weight=1.0)
        text(s, x, 3.62, w, 0.5, [r(cn, size=22, bold=True, color=INK, latin=F_BLACK)])
        text(s, x, 4.1, w, 0.3, [r(en, size=10, color=RED, spc=200, latin=F_SANS)])
        text(s, x, 4.5, w, 1.3, [r(desc, size=13.5, color=GRAY, latin=F_SANS, ea=F_EA)], ls=1.35)
    rect(s, M, 6.25, CW, 0.55, fill=INK)
    text(s, M + 0.2, 6.25, CW - 0.4, 0.55,
         [r("用 AI Agent 替代「监测员 + 分析师 + 合规审核」组合 ——", size=12.5, color=WHITE, latin=F_SANS, ea=F_EA),
          r("  从「看见」到「看懂、研判、行动」的闭环。", size=12.5, bold=True, color=WHITE, latin=F_SANS, ea=F_EA)],
         anchor=MSO_ANCHOR.MIDDLE)
    notes(s, "我们用 AI Agent 替代监测员、分析师、合规审核的组合,实现从看见到看懂、研判、行动的闭环。这是资深舆情分析师的完整工作流,首次被工程化。")


def slide_06():
    s = add_slide()
    page_head(s, 6, "量化价值", "Validated Outcomes · City Commercial Bank POC")
    kpis = [
        ("60%", "误报率下降", "FALSE-POSITIVE  vs 通用工具"),
        ("10min", "负面平均响应", "4h  →  10min"),
        ("×3", "分析师产能", "单日覆盖事件数"),
        ("92%+", "实体识别准确率", "金融 NER 准确率"),
    ]
    c = cols(4, gap=0.5)
    for i, (big, cn, en) in enumerate(kpis):
        x, w = c[i]
        text(s, x, 2.7, w, 1.2, [r(big, size=58, bold=True, color=INK, latin=F_BLACK)], ls=1.0)
        hline(s, x, 3.95, w, color=RED, weight=2.5)
        text(s, x, 4.1, w, 0.4, [r(cn, size=15, bold=True, color=INK, latin=F_BLACK)])
        text(s, x, 4.5, w, 0.3, [r(en, size=9.5, color=GRAY, spc=120, latin=F_SANS)])
    rect(s, M, 5.5, CW, 0.5, fill=BG2)
    text(s, M + 0.15, 5.5, CW - 0.3, 0.5,
         [r("数据来源:某城商行消保舆情 POC 验证  ·  标 [核实] 处需替换为真实 POC 数据", size=11, color=GRAY, latin=F_SANS, ea=F_EA)],
         anchor=MSO_ANCHOR.MIDDLE)
    notes(s, "这是我们在某城商行消保舆情场景的 POC 验证结果。误报率下降六成,响应时长从小时级压缩到十分钟级,分析师产能提升三倍。这些数字需在路演前替换为真实 POC 数据。")


def slide_07():
    s = add_slide()
    page_head(s, 7, "为什么是现在", "Why Now · Timing")
    fac = [
        ("技术", "TECH", "大模型 + Agent 框架成熟,「理解 + 行动」首次工程可行。"),
        ("监管", "POLICY", "声誉风险、消保、投资者保护、ESG、数据合规要求趋严。"),
        ("基建", "INFRA", "金融机构数据治理就绪,私有化部署大模型成为可能。"),
        ("市场", "MARKET", "通用舆情存量竞争,「智能化升级」窗口打开。"),
    ]
    c = cols(4, gap=0.5)
    for i, (cn, en, desc) in enumerate(fac):
        x, w = c[i]
        text(s, x, 2.6, w, 0.4, [r("0%d" % (i + 1), size=12, color=RED, spc=200, latin=F_SANS)])
        text(s, x, 2.95, w, 0.7, [r(cn, size=30, bold=True, color=INK, latin=F_BLACK)])
        text(s, x, 3.7, w, 0.3, [r(en, size=10, color=GRAY, spc=200, latin=F_SANS)])
        hline(s, x, 4.05, w, color=INK, weight=1.0)
        text(s, x, 4.25, w, 1.6, [r(desc, size=13.5, color=GRAY, latin=F_SANS, ea=F_EA)], ls=1.4)
    notes(s, "三个趋势同时成熟:模型能力到了,监管要求到了,客户基建也到了。这是金融舆情智能体最好的时间窗。")


def slide_09():
    s = add_slide()
    page_head(s, 9, "产品架构", "Layered Architecture")
    layers = [
        ("应用层", "APPLICATION", "舆情驾驶舱 · 研判报告 · 实时预警 · 投关助手 · 危机指挥"),
        ("Agent 编排层", "ORCHESTRATION", "监测 / 研判 / 应对 / 合规 / 预警 多 Agent 协同 · 任务规划与工具调用"),
        ("金融知识图谱层", "KNOWLEDGE GRAPH", "公司—人物—事件—关联交易关系网 · 支撑关联研判与溯源"),
        ("智能处理层", "NLP ENGINE", "金融 NER · 情感分析 · 事件抽取 · 谣言检测 · 跨语言"),
        ("数据采集层", "DATA INTAKE", "全网爬虫 + 授权数据源 + API · 社媒/新闻/股吧/公告/研报/境外媒体"),
        ("模型层", "MODEL", "金融行业大模型 + 私有化部署 + 可插拔主流模型"),
    ]
    y0 = 2.45; h = 0.62; gap = 0.12
    for i, (cn, en, desc) in enumerate(layers):
        y = y0 + i * (h + gap)
        alt = (i % 2 == 1)
        rect(s, M, y, CW, h, fill=(BG2 if alt else WHITE), line=BORDER, weight=0.75)
        rect(s, M, y, 0.09, h, fill=RED)
        text(s, M + 0.32, y, 3.0, h, [r(cn, size=15, bold=True, color=INK, latin=F_BLACK)], anchor=MSO_ANCHOR.MIDDLE)
        text(s, M + 3.3, y, 2.6, h, [r(en, size=9.5, color=GRAY, spc=160, latin=F_SANS)], anchor=MSO_ANCHOR.MIDDLE)
        text(s, M + 5.9, y, CW - 6.1, h, [r(desc, size=11.5, color=GRAY, latin=F_SANS, ea=F_EA)], anchor=MSO_ANCHOR.MIDDLE)
    notes(s, "六层架构自下而上:数据采集、智能处理、金融知识图谱、Agent 编排、应用,以及底座的金融行业大模型。每一层都为「金融语义 + 可解释 + 私有化」服务。")


def slide_10():
    s = add_slide()
    page_head(s, 10, "核心技术与差异化", "Core Tech & Differentiation")
    q = [
        ("01", "金融语义理解", "FINANCIAL SEMANTICS", "领域微调 + 知识图谱,远超通用舆情的关键词命中。"),
        ("02", "秒级实时引擎", "REAL-TIME ENGINE", "流式处理,潜在风险分钟 / 秒级预警。"),
        ("03", "幻觉抑制 + 可解释", "HALLUCINATION CONTROL", "结论附证据链与出处,金融场景必备。"),
        ("04", "多 Agent + 低代码", "MULTI-AGENT & LOW-CODE", "业务人员可自定义研判规则与 Agent。"),
    ]
    c = cols(2, gap=0.7); ys = [2.5, 4.62]
    for i, (no, cn, en, desc) in enumerate(q):
        x, w = c[i % 2]; y = ys[i // 2]
        text(s, x, y, w, 0.4, [r(no, size=12, color=RED, spc=200, latin=F_SANS)])
        text(s, x, y + 0.32, w, 0.55, [r(cn, size=22, bold=True, color=INK, latin=F_BLACK)])
        text(s, x, y + 0.9, w, 0.3, [r(en, size=10, color=GRAY, spc=200, latin=F_SANS)])
        hline(s, x, y + 1.25, w, color=INK, weight=1.0)
        text(s, x, y + 1.4, w, 0.6, [r(desc, size=13, color=GRAY, latin=F_SANS, ea=F_EA)], ls=1.3)
    notes(s, "四个差异化点:金融语义、秒级实时、幻觉抑制加可解释、多 Agent 加低代码配置。这是证明我们不是套壳的关键。")


def slide_11():
    s = add_slide()
    page_head(s, 11, "部署形态与安全合规", "Deployment & Compliance")
    modes = [
        ("私有化", "PRIVATE", "银行 / 监管首选", "全栈本地化 · 数据不出域"),
        ("混合云", "HYBRID", "中大型机构", "模型与敏感数据本地 · 弹性算力上云"),
        ("SaaS", "SAAS", "中小机构 / 上市公司", "快速接入 · 按账号与数据量计费"),
    ]
    c = cols(3, gap=0.6)
    for i, (cn, en, who, desc) in enumerate(modes):
        x, w = c[i]
        rect(s, x, 2.55, w, 2.85, fill=WHITE, line=BORDER, weight=1.0)
        rect(s, x, 2.55, w, 0.1, fill=RED)
        text(s, x + 0.25, 2.85, w - 0.5, 0.6, [r(cn, size=26, bold=True, color=INK, latin=F_BLACK)])
        text(s, x + 0.25, 3.5, w - 0.5, 0.3, [r(en, size=10, color=GRAY, spc=200, latin=F_SANS)])
        hline(s, x + 0.25, 3.85, w - 0.5, color=INK, weight=1.0)
        text(s, x + 0.25, 4.0, w - 0.5, 0.4, [r(who, size=13, bold=True, color=INK, latin=F_SANS, ea=F_EA)])
        text(s, x + 0.25, 4.4, w - 0.5, 0.8, [r(desc, size=12, color=GRAY, latin=F_SANS, ea=F_EA)], ls=1.3)
    rect(s, M, 5.75, CW, 0.6, fill=INK)
    text(s, M + 0.2, 5.75, CW - 0.4, 0.6,
         [r("合规硬门槛  ·  ", size=11, bold=True, color=RED, latin=F_SANS, ea=F_EA),
          r("权限分级  ·  操作审计  ·  数据脱敏  ·  信创适配  ·  等保 [核实]", size=11, color=WHITE, latin=F_SANS, ea=F_EA)],
         anchor=MSO_ANCHOR.MIDDLE)
    notes(s, "三档部署形态覆盖不同客户:银行监管走私有化,中大型机构走混合云,中小机构走 SaaS。权限分级、操作审计、脱敏、信创适配和等保是 To B / To G 的硬门槛。")


def slide_12():
    s = add_slide()
    page_head(s, 12, "产品路线图", "Roadmap")
    stops = [
        ("Q1–Q2", "平台 GA", "核心平台上线 + 金融知识图谱 v1"),
        ("Q3–Q4", "私有化闭环", "私有化部署闭环 + 3 个标杆客户"),
        ("Y2", "行业版扩展", "消保 / 声誉 / 投关 / 地方金融风险"),
    ]
    y = 4.2
    hline(s, M + 0.4, y, CW - 0.8, color=INK, weight=1.5)
    c = cols(3, gap=0.6)
    for i, (per, cn, desc) in enumerate(stops):
        x, w = c[i]
        rect(s, x + 0.2, y - 0.09, 0.18, 0.18, fill=RED)
        text(s, x, y - 1.15, w, 0.4, [r(per, size=13, bold=True, color=RED, spc=160, latin=F_SANS)])
        text(s, x, y - 0.7, w, 0.5, [r(cn, size=20, bold=True, color=INK, latin=F_BLACK)])
        text(s, x, y + 0.35, w, 1.0, [r(desc, size=13, color=GRAY, latin=F_SANS, ea=F_EA)], ls=1.4)
    notes(s, "路线图:前两季度平台 GA 加知识图谱 v1;下半年跑通私有化闭环并拿下三个标杆;第二年开始行业版扩展。")


def slide_14():
    s = add_slide()
    page_head(s, 14, "市场规模", "TAM · SAM · SOM")
    rows = [
        ("TAM", "总量市场", "企业舆情监测与公关分析服务市场", "数十亿级 RMB", "[艾瑞 / 赛迪 / IDC 核实]"),
        ("SAM", "可服务市场", "金融垂直:声誉 / 消保 / 投关支出", "数亿级 RMB", "银行 / 券商 / 上市公司 / 监管"),
        ("SOM", "3 年可获取", "目标客户数 × 渗透率 × 客单价", "数千万级 RMB", "私有化百万级 · SaaS 数十万/年"),
    ]
    y0 = 2.55; h = 1.15; gap = 0.18
    for i, (code, cn, desc, scale, note) in enumerate(rows):
        y = y0 + i * (h + gap)
        rect(s, M, y, CW, h, fill=(BG2 if i % 2 else WHITE), line=BORDER, weight=0.75)
        text(s, M + 0.3, y, 1.6, h, [r(code, size=30, bold=True, color=RED, latin=F_BLACK)], anchor=MSO_ANCHOR.MIDDLE)
        vline(s, M + 1.95, y + 0.18, h - 0.36, color=BORDER, weight=1.0)
        text(s, M + 2.25, y, 4.2, h, [r(cn, size=16, bold=True, color=INK, latin=F_BLACK)], anchor=MSO_ANCHOR.MIDDLE)
        text(s, M + 2.25, y + 0.45, 4.2, h, [r(desc, size=12, color=GRAY, latin=F_SANS, ea=F_EA)], anchor=MSO_ANCHOR.TOP)
        text(s, M + 6.8, y, 3.4, h, [r(scale, size=22, bold=True, color=INK, latin=F_BLACK)], anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.RIGHT)
        text(s, M + 6.8, y + 0.5, 3.4, h, [r(note, size=10, color=LGRAY, latin=F_SANS, ea=F_EA)], anchor=MSO_ANCHOR.TOP, align=PP_ALIGN.RIGHT)
    notes(s, "市场规模三层:TAM 是企业舆情监测与公关分析服务市场,数十亿级;SAM 收敛到金融垂直支出;SOM 用自下而上口径——目标客户数乘渗透率乘客单价测算,3 年可获取数千万级。所有数字需权威报告核实。")


def slide_15():
    s = add_slide()
    page_head(s, 15, "竞争定位矩阵", "Competitive Positioning")
    bx, by, bw, bh = M, 2.55, CW, 3.7
    rect(s, bx, by, bw, bh, fill=WHITE, line=BORDER, weight=1.0)
    hline(s, bx, by + bh / 2, bw, color=BORDER, weight=1.0)
    vline(s, bx + bw / 2, by, bh, color=BORDER, weight=1.0)
    text(s, bx, by - 0.32, bw, 0.3, [r("Agent 智能  ↑", size=10, color=GRAY, spc=120, latin=F_SANS)], align=PP_ALIGN.LEFT)
    text(s, bx, by - 0.32, bw, 0.3, [r("规则驱动", size=10, color=LGRAY, latin=F_SANS)], align=PP_ALIGN.RIGHT)
    text(s, bx, by + bh + 0.08, bw, 0.3, [r("←  通用", size=10, color=GRAY, latin=F_SANS, ea=F_EA)], align=PP_ALIGN.LEFT)
    text(s, bx, by + bh + 0.08, bw, 0.3, [r("金融垂直  →", size=10, color=GRAY, latin=F_SANS, ea=F_EA)], align=PP_ALIGN.RIGHT)
    dots = [
        (0.18, 0.74, "传统舆情厂商", INK),
        (0.62, 0.30, "通用大模型平台", INK),
        (0.20, 0.30, "海外财经情报", INK),
        (0.60, 0.55, "金融机构自建 IT", INK),
        (0.84, 0.20, "天微 AI 实验室", RED),
    ]
    for rx, ry, label, col in dots:
        px = bx + rx * bw; py = by + ry * bh
        rect(s, px - 0.07, py - 0.07, 0.14, 0.14, fill=col)
        text(s, px + 0.14, py - 0.16, 3.0, 0.32,
             [r(label, size=11.5, bold=(col == RED), color=col, latin=F_SANS, ea=F_EA)])
    text(s, bx + bw - 3.4, by + bh - 0.5, 3.2, 0.4,
         [r("空白生态位 = 金融垂直 × Agent 智能 × 私有化合规", size=10.5, color=RED, latin=F_SANS, ea=F_EA)],
         align=PP_ALIGN.RIGHT)
    notes(s, "四类对手各有短板:传统舆情规则强但智能弱;通用大模型平台私有化弱;海外财经情报本地合规弱;自建 IT 成本高。我们的位置是金融垂直、Agent 智能、私有化合规三者的交叉点,目前没人占住。")


def slide_16():
    s = add_slide()
    page_head(s, 16, "目标客户与场景", "Target Customers & Scenarios")
    cards = [
        ("01", "城商行 / 农商行", "消保舆情 · 品牌声誉"),
        ("02", "券商", "研报合规 · 声誉管理"),
        ("03", "上市公司董办", "投资者关系 · 危机公关"),
        ("04", "地方金融监管 (To G)", "非法集资监测 · 地方金融风险"),
    ]
    c = cols(2, gap=0.7); ys = [2.5, 4.4]
    for i, (no, cn, desc) in enumerate(cards):
        x, w = c[i % 2]; y = ys[i // 2]
        rect(s, x, y, w, 1.65, fill=WHITE, line=BORDER, weight=1.0)
        rect(s, x, y, 0.09, 1.65, fill=RED)
        text(s, x + 0.3, y + 0.18, w - 0.5, 0.4, [r(no, size=12, color=RED, spc=200, latin=F_SANS)])
        text(s, x + 0.3, y + 0.5, w - 0.5, 0.5, [r(cn, size=19, bold=True, color=INK, latin=F_BLACK)])
        text(s, x + 0.3, y + 1.05, w - 0.5, 0.5, [r(desc, size=12.5, color=GRAY, latin=F_SANS, ea=F_EA)])
    rect(s, M, 6.3, CW, 0.5, fill=BG2)
    text(s, M + 0.2, 6.3, CW - 0.4, 0.5,
         [r("决策链  ", size=11, bold=True, color=INK, latin=F_SANS, ea=F_EA),
          r("品牌 / 合规 / 办公室 / IT      ", size=11, color=GRAY, latin=F_SANS, ea=F_EA),
          r("采购周期  ", size=11, bold=True, color=INK, latin=F_SANS, ea=F_EA),
          r("3–9 个月      ", size=11, color=GRAY, latin=F_SANS, ea=F_EA),
          r("预算来源  ", size=11, bold=True, color=INK, latin=F_SANS, ea=F_EA),
          r("科技 / 品牌 / 合规", size=11, color=GRAY, latin=F_SANS, ea=F_EA)],
         anchor=MSO_ANCHOR.MIDDLE)
    notes(s, "首发客户是城商行和上市公司董办,决策相对快、预算明确;中期延伸到券商和地方金融监管。")


def slide_17():
    s = add_slide()
    page_head(s, 17, "竞品对比", "Competitive Comparison")
    headers = ["维度", "天微", "通用舆情", "通用大模型", "人工"]
    data = [
        ["金融语义理解", "★★★★★", "★★", "★★★", "★★★★"],
        ["响应速度",     "秒级",   "分钟级", "分钟级", "小时级"],
        ["可解释 / 留痕", "强",     "中",   "弱",    "强"],
        ["私有化合规",   "原生支持", "弱",   "弱",    "—"],
        ["成本",         "中",     "低",   "中",    "高"],
    ]
    x0 = M; y0 = 2.5; rh = 0.62
    widths = [3.0, 2.2, 2.2, 2.2, 2.39]
    cx = x0
    for j, htxt in enumerate(headers):
        fill = INK if j == 1 else BG2
        rect(s, cx, y0, widths[j], rh, fill=fill, line=BORDER, weight=0.75)
        colr = WHITE if j == 1 else INK
        text(s, cx, y0, widths[j], rh,
             [r(htxt, size=12.5, bold=True, color=colr, latin=F_BLACK, ea=F_EA)],
             anchor=MSO_ANCHOR.MIDDLE, align=(PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER))
        cx += widths[j]
    for i, row in enumerate(data):
        cy = y0 + (i + 1) * rh
        cx = x0
        for j, val in enumerate(row):
            fill = WHITE
            if j == 1:
                fill = SOFTRED
            elif i % 2 == 1:
                fill = BG2
            rect(s, cx, cy, widths[j], rh, fill=fill, line=BORDER, weight=0.75)
            colr = RED if j == 1 else (INK if j == 0 else GRAY)
            text(s, cx, cy, widths[j], rh,
                 [r(val, size=12.5, bold=(j == 1 or j == 0), color=colr, latin=F_SANS, ea=F_EA)],
                 anchor=MSO_ANCHOR.MIDDLE, align=(PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER))
            cx += widths[j]
    notes(s, "横向对比可以看到,我们是唯一在金融语义、可解释、私有化三项同时达标的方案。")


def slide_19():
    s = add_slide()
    page_head(s, 19, "商业模式", "Revenue Model · Four Tiers")
    tiers = [
        ("01", "订阅 SaaS", "SUBSCRIPTION", "按账号 / 模块 / 数据量计费", "中小机构 · 上市公司"),
        ("02", "私有化项目", "ON-PREMISE", "一次性交付 + 年度维保", "银行 · 监管"),
        ("03", "按量计费", "USAGE-BASED", "API 调用 / Token 计费", "集成方 · ISV"),
        ("04", "行业解决方案", "SOLUTION", "打包交付 · 场景化定价", "消保 / 声誉 / 投关 / 地方金融风险"),
    ]
    c = cols(4, gap=0.45)
    for i, (no, cn, en, how, who) in enumerate(tiers):
        x, w = c[i]
        rect(s, x, 2.55, w, 3.4, fill=WHITE, line=BORDER, weight=1.0)
        rect(s, x, 2.55, w, 0.1, fill=RED)
        text(s, x + 0.22, 2.8, w - 0.4, 0.4, [r(no, size=12, color=RED, spc=200, latin=F_SANS)])
        text(s, x + 0.22, 3.15, w - 0.4, 0.6, [r(cn, size=17, bold=True, color=INK, latin=F_BLACK)], ls=1.0)
        text(s, x + 0.22, 3.78, w - 0.4, 0.3, [r(en, size=9, color=GRAY, spc=180, latin=F_SANS)])
        hline(s, x + 0.22, 4.15, w - 0.44, color=INK, weight=1.0)
        text(s, x + 0.22, 4.3, w - 0.4, 0.8, [r(how, size=12, color=GRAY, latin=F_SANS, ea=F_EA)], ls=1.3)
        text(s, x + 0.22, 5.35, w - 0.4, 0.5, [r(who, size=11, bold=True, color=INK, latin=F_SANS, ea=F_EA)], ls=1.2)
    notes(s, "四档商业模式覆盖不同客户:中小走 SaaS,银行走私有化,生态伙伴走按量计费,行业客户走解决方案。收入结构会随客户成熟度逐步转向高毛利的订阅与私有化维保。")


def slide_20():
    s = add_slide()
    page_head(s, 20, "18–24 月里程碑", "Milestones")
    ms = [
        ("产品", "PRODUCT", "平台 GA + 金融知识图谱 v1 + 私有化部署闭环"),
        ("客户", "CUSTOMER", "3 个标杆付费客户(城商行 / 券商 / 上市公司)"),
        ("能力", "CAPABILITY", "实时预警秒级 + 可解释证据链 + 信创适配"),
        ("融资", "FUNDING", "种子轮 1,000 万 → 18–24 月后启动 Pre-A"),
    ]
    c = cols(4, gap=0.5)
    for i, (cn, en, desc) in enumerate(ms):
        x, w = c[i]
        text(s, x, 2.6, w, 0.4, [r("0%d" % (i + 1), size=12, color=RED, spc=200, latin=F_SANS)])
        text(s, x, 2.95, w, 0.6, [r(cn, size=26, bold=True, color=INK, latin=F_BLACK)])
        text(s, x, 3.65, w, 0.3, [r(en, size=10, color=GRAY, spc=200, latin=F_SANS)])
        hline(s, x, 4.0, w, color=INK, weight=1.0)
        text(s, x, 4.2, w, 2.0, [r(desc, size=13, color=GRAY, latin=F_SANS, ea=F_EA)], ls=1.5)
    notes(s, "这四个里程碑是种子轮的核心交付:产品、客户、能力、下一轮。达成后我们就有 Pre-A 的数据支撑。")


def slide_21():
    s = add_slide()
    page_head(s, 21, "资金用途", "Use of Proceeds · Seed Round")
    text(s, M, 2.55, 4.5, 1.0, [r("1,000", size=72, bold=True, color=INK, latin=F_BLACK)], ls=1.0)
    text(s, M, 3.7, 4.5, 0.4, [r("万 RMB  ·  跑道 18–24 个月", size=15, color=GRAY, latin=F_SANS, ea=F_EA)])
    hline(s, M, 4.25, 4.3, color=RED, weight=2.5)
    text(s, M, 4.4, 4.5, 0.4, [r("种子轮重压研发与场景验证", size=12, color=INK, latin=F_SANS, ea=F_EA)])
    seg = [("研发与算力", 50, "500 万", INK),
           ("金融数据与知识图谱", 18, "180 万", RED),
           ("标杆客户与场景", 15, "150 万", BLUE),
           ("团队建设", 12, "120 万", GRAY),
           ("运营 / 合规 / 储备", 5, "50 万", BORDER)]
    # 100% 堆叠条 · 瑞士式平涂(替代原生环形图)
    bar_x, bar_y, bar_w, bar_h = M, 4.95, CW, 0.72
    cx2 = bar_x
    for nm, pct, amt, col in seg:
        sw2 = bar_w * pct / 100.0
        rect(s, cx2, bar_y, sw2, bar_h, fill=col)
        if pct >= 12:
            lblc = WHITE if col != BORDER else INK
            text(s, cx2, bar_y, sw2, bar_h, [r("%d%%" % pct, size=14, bold=True, color=lblc, latin=F_SANS)],
                 anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)
        cx2 += sw2
    cx2 = bar_x
    for nm, pct, amt, col in seg:
        sw2 = bar_w * pct / 100.0
        text(s, cx2 + 0.05, bar_y + bar_h + 0.12, sw2 - 0.1, 0.3, [r(nm, size=10.5, color=INK, latin=F_SANS, ea=F_EA)], ls=1.0)
        text(s, cx2 + 0.05, bar_y + bar_h + 0.42, sw2 - 0.1, 0.3, [r(amt, size=10, color=GRAY, latin=F_SANS, ea=F_EA)], ls=1.0)
        cx2 += sw2
    notes(s, "1000 万的分配逻辑:种子轮重压研发与场景验证。研发算力五成,数据与知识图谱近两成,客户与场景一成半,团队一成出头,运营合规留个尾数。跑道 18 到 24 个月。")


def slide_22():
    s = add_slide()
    page_num(s, 22)
    text(s, M, 0.95, 11, 0.5, [r("CLOSING", size=12, color=GRAY, spc=260, latin=F_SANS)])
    text(s, M, 1.4, 12, 1.6,
         [r("验证 PMF,", size=44, bold=True, color=INK, latin=F_BLACK, ea=F_EA),
          r("  拿下标杆,", size=44, bold=True, color=INK, latin=F_BLACK, ea=F_EA)], ls=1.1)
    text(s, M, 2.45, 12, 1.0,
         [r("跑通私有化闭环。", size=44, bold=True, color=RED, latin=F_BLACK, ea=F_EA)], ls=1.1)
    hline(s, M, 3.6, 0.72, color=RED, weight=3.0)
    c = cols(3, gap=0.6)
    blocks = [
        ("本轮", "1,000 万 RMB · 种子轮 · 跑道 18–24 月"),
        ("估值", "投后约 5,000 万–1 亿(出让 10–20%)[内部决策]"),
        ("退出路径", "下一轮 Pre-A / 战略并购(金融科技 · 舆情 · 数据厂商)"),
    ]
    for i, (k, v) in enumerate(blocks):
        x, w = c[i]
        text(s, x, 3.85, w, 0.35, [r(k.upper(), size=11, color=RED, spc=220, latin=F_SANS, ea=F_EA)])
        hline(s, x, 4.2, w, color=INK, weight=1.0)
        text(s, x, 4.35, w, 1.4, [r(v, size=13.5, color=INK, latin=F_SANS, ea=F_EA)], ls=1.4)
    rect(s, M, 6.0, CW, 0.55, fill=INK)
    text(s, M + 0.2, 6.0, CW - 0.4, 0.55,
         [r("核心团队  ·  ", size=11, bold=True, color=RED, latin=F_SANS, ea=F_EA),
          r("AI 工程 + 金融行业 + To B 销售  [贵司提供履历 / 照片]", size=11, color=WHITE, latin=F_SANS, ea=F_EA)],
         anchor=MSO_ANCHOR.MIDDLE)
    page_foot(s, 22)
    notes(s, "总结:我们用 1000 万种子轮验证金融舆情智能体的 PMF,拿三个标杆客户,跑通私有化闭环。下一轮 Pre-A 时,我们希望带着真实的付费客户和单位经济数据,和各位再深入交流。谢谢!")


SLIDES = [slide_01, slide_02, slide_03, slide_04, slide_05, slide_06, slide_07,
          slide_08, slide_09, slide_10, slide_11, slide_12, slide_13, slide_14,
          slide_15, slide_16, slide_17, slide_18, slide_19, slide_20, slide_21, slide_22]

for fn in SLIDES:
    fn()

OUT = r"G:\cz-night-learn-20260713\GLM\002\PPT\商业计划书-瑞士风-天微AI实验室.pptx"
prs.save(OUT)
print("SAVED:", OUT)
print("SLIDES:", len(prs.slides._sldIdLst))
